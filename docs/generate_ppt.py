#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
义乌小商品出海智能体-OPC 路演PPT生成脚本 V2冠军版
使用 python-pptx 生成专业路演演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 颜色定义 ====================
YIWU_RED = RGBColor(0xD4, 0x27, 0x2C)       # 义乌红主色
DARK_RED = RGBColor(0xA0, 0x1E, 0x22)        # 深红
LIGHT_RED = RGBColor(0xE8, 0x5C, 0x60)       # 浅红
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GOLD = RGBColor(0xD4, 0xA0, 0x17)            # 金色点缀
ACCENT_BLUE = RGBColor(0x2C, 0x5F, 0x8A)     # 辅助蓝
ACCENT_GREEN = RGBColor(0x2D, 0x8B, 0x4E)    # 辅助绿

# ==================== PPT尺寸 ====================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

TOTAL_PAGES = 22


# ==================== 工具函数 ====================
def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=DARK_GRAY,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', line_spacing=1.2):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16, font_color=DARK_GRAY,
                   bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', line_spacing=1.5):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(2)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, font_color=DARK_GRAY,
                    bullet_char='●', font_name='微软雅黑', line_spacing=1.6, indent=Pt(20)):
    """添加带项目符号的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet_char}  {item}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_page_header(slide, title, subtitle=None):
    """添加页面顶部标题栏"""
    # 顶部红色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), fill_color=YIWU_RED)
    # 标题
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7), title,
                font_size=32, font_color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4), subtitle,
                    font_size=16, font_color=RGBColor(0xFF, 0xCC, 0xCC), bold=False, alignment=PP_ALIGN.LEFT)
    # 底部细线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.04), fill_color=GOLD)


def add_page_number(slide, num):
    """添加页码"""
    add_textbox(slide, Inches(12), Inches(7.05), Inches(1.2), Inches(0.35),
                f'{num} / {TOTAL_PAGES}', font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, content_lines, title_color=YIWU_RED,
             bg_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD)):
    """添加卡片式内容块"""
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill_color=bg_color, line_color=border_color, line_width=Pt(1))
    add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
              width - Inches(0.1), Inches(0.06), fill_color=title_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.4),
                title, font_size=18, font_color=title_color, bold=True)
    y_offset = top + Inches(0.65)
    for line in content_lines:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                    line, font_size=13, font_color=DARK_GRAY)
        y_offset += Inches(0.3)
    return card


# ==================== 第1页：封面页 ====================
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 全屏红色背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=YIWU_RED)
    # 装饰：右侧深色块
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0), Inches(3.833), SLIDE_HEIGHT, fill_color=DARK_RED)
    # 装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(2), Inches(0.06), fill_color=GOLD)
    # V2冠军版标签
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(2.2), Inches(0.5),
              fill_color=GOLD, line_color=None)
    add_textbox(slide, Inches(0.8), Inches(0.63), Inches(2.2), Inches(0.45),
                'V2冠军版', font_size=22, font_color=DARK_RED, bold=True, alignment=PP_ALIGN.CENTER)
    # 项目名称
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.2),
                '义乌小商品出海智能体', font_size=48, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.9), Inches(8.5), Inches(0.8),
                'OPC', font_size=72, font_color=GOLD, bold=True)
    # Slogan
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(8.5), Inches(0.6),
                '让小商品照亮全球', font_size=28, font_color=RGBColor(0xFF, 0xDD, 0xDD), bold=False)
    # 路演人信息
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(3), Inches(0.02), fill_color=RGBColor(0xFF, 0xAA, 0xAA))
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(5), Inches(0.4),
                '路演人：冯亦根', font_size=20, font_color=WHITE, bold=False)
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.4),
                'Yiwu Internet Competition 2026', font_size=14, font_color=RGBColor(0xFF, 0xBB, 0xBB))
    add_page_number(slide, 1)


# ==================== 第2页：痛点页 ====================
def slide_pain_points():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '出海之痛', '义乌商户出海面临的五大核心痛点')

    pain_items = [
        ('语', '语言不通', '多语言沟通障碍，无法有效触达海外客户'),
        ('渠', '渠道不熟', '缺乏海外销售渠道，依赖中间商利润被压缩'),
        ('规', '合规不清', '各国法规差异大，合规风险高成本高'),
        ('运', '运营不精', '跨境电商运营门槛高，人才稀缺'),
        ('物', '物流不懂', '国际物流链路复杂，时效成本难控'),
    ]

    for i, (icon, title, desc) in enumerate(pain_items):
        col = i % 5
        left = Inches(0.6 + col * 2.5)
        top = Inches(2.0)

        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(4.5),
                         fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(2.1), Inches(0.06), fill_color=YIWU_RED)
        circle_size = Inches(0.8)
        add_shape(slide, MSO_SHAPE.OVAL,
                  left + (Inches(2.2) - circle_size) / 2, top + Inches(0.4),
                  circle_size, circle_size, fill_color=YIWU_RED)
        add_textbox(slide, left + (Inches(2.2) - circle_size) / 2, top + Inches(0.48),
                    circle_size, circle_size, f'0{i+1}',
                    font_size=22, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(1.5), Inches(2.0), Inches(0.4),
                    title, font_size=22, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.6), top + Inches(2.1),
                  Inches(1.0), Inches(0.03), fill_color=GOLD)
        add_textbox(slide, left + Inches(0.15), top + Inches(2.4), Inches(1.9), Inches(1.8),
                    desc, font_size=13, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 2)


# ==================== 第3页：市场规模页 ====================
def slide_market_size():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '市场机遇', '全球小商品市场万亿级蓝海')

    data_items = [
        ('$5000亿', '全球小商品市场规模', '年增长率8.5%\n跨境电商增速超20%'),
        ('¥5000亿', '义乌年出口额', '占全国小商品出口1/3\n通达230+国家和地区'),
        ('7.5万', '义乌国际商贸城商铺', '210万SKU\n全球最大小商品集散中心'),
    ]

    for i, (num, title, desc) in enumerate(data_items):
        left = Inches(0.8 + i * 4.1)
        top = Inches(2.0)

        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(4.5),
                         fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(3.5), Inches(0.08), fill_color=YIWU_RED)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.5), Inches(3.2), Inches(1.2),
                    num, font_size=52, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.8), Inches(3.2), Inches(0.5),
                    title, font_size=20, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(1.0), top + Inches(2.5),
                  Inches(1.6), Inches(0.03), fill_color=GOLD)
        for j, line in enumerate(desc.split('\n')):
            add_textbox(slide, left + Inches(0.3), top + Inches(2.8 + j * 0.5), Inches(3.0), Inches(0.4),
                        line, font_size=15, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 3)


# ==================== 第4页：义乌发展经验国家战略页（新增） ====================
def slide_yiwu_strategy():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '"义乌发展经验"国家战略', '习近平总书记亲自批示，国家级制度性壁垒')

    # 顶部战略标识
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.7),
              fill_color=RGBColor(0xFF, 0xF0, 0xF0), line_color=YIWU_RED, line_width=Pt(2))
    add_textbox(slide, Inches(0.8), Inches(1.68), Inches(11.7), Inches(0.5),
                '★ 习近平总书记亲自批示 ——"义乌发展经验"是国家级改革开放战略，不可复制的制度性壁垒',
                font_size=16, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)

    # 左侧：义乌发展经验核心
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.6), Inches(5.8), Inches(4.5),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(0.9), Inches(2.8), Inches(5.2), Inches(0.5),
                '义乌发展经验 · 六个坚持', font_size=22, font_color=WHITE, bold=True)

    six_principles = [
        '坚持兴商建市——以市场为核心驱动力',
        '坚持产业联动——贸工联动协调发展',
        '坚持统筹城乡——城乡一体化发展',
        '坚持和谐发展——经济与社会协调',
        '坚持丰厚文化——文化软实力支撑',
        '坚持党政有为——有为政府与有效市场',
    ]
    for i, item in enumerate(six_principles):
        y = Inches(3.5 + i * 0.55)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(5.2), Inches(0.45),
                  fill_color=DARK_RED, line_color=None)
        add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(4.8), Inches(0.35),
                    item, font_size=14, font_color=WHITE, bold=False)

    # 右侧：国家战略价值
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.5),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.8), Inches(5.2), Inches(0.5),
                '国家战略价值', font_size=22, font_color=YIWU_RED, bold=True)

    strategy_items = [
        ('1039模式', '义乌首创市场采购贸易方式，全国39城复制推广'),
        ('210万企业', '义乌模式带动全国210万市场主体发展'),
        ('3200万就业', '义乌经验支撑3200万人就业'),
        ('政策网络', '39城政策互通，形成全国性制度网络'),
        ('不可复制', '地理基因+首创经验+政策网络，三位一体壁垒'),
    ]
    for i, (title, desc) in enumerate(strategy_items):
        y = Inches(3.5 + i * 0.7)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), y, Inches(5.2), Inches(0.6),
                  fill_color=LIGHT_GRAY, line_color=None)
        add_textbox(slide, Inches(7.3), y + Inches(0.05), Inches(1.8), Inches(0.3),
                    title, font_size=14, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(9.2), y + Inches(0.05), Inches(2.9), Inches(0.5),
                    desc, font_size=12, font_color=MID_GRAY)

    add_page_number(slide, 4)


# ==================== 第5页：解决方案总览页（7个Agent） ====================
def slide_solution_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '解决方案总览', '7大AI Agent 全链路覆盖出海全流程')

    agents = [
        ('1', '市场洞察\nAgent', '实时市场趋势\n义乌指数分析'),
        ('2', '智能选品\nAgent', '210万SKU\n智能推荐'),
        ('3', '供应链匹配\nAgent', '7.5万商铺\n精准对接'),
        ('4', '跨境内容\nAgent', '8语言4平台\n一键生成'),
        ('5', '合规助手\nAgent', '1039模式\n沿线合规'),
        ('6', '智能客服\nAgent', '7×24小时\n多语言服务'),
        ('7', '政策复制\nAgent', '39城政策\n一键复制'),
    ]

    for i, (num, title, desc) in enumerate(agents):
        left = Inches(0.3 + i * 1.85)
        top = Inches(2.2)

        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.7), Inches(4.0),
                         fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        # 顶部色块
        color = YIWU_RED if i < 6 else GOLD
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(1.6), Inches(0.8), fill_color=color)
        # 编号
        add_textbox(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.6), Inches(0.7),
                    num, font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Agent名称
        for j, line in enumerate(title.split('\n')):
            add_textbox(slide, left + Inches(0.05), top + Inches(1.1 + j * 0.4), Inches(1.6), Inches(0.35),
                        line, font_size=14, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
        # 分隔线
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.3), top + Inches(2.1),
                  Inches(1.1), Inches(0.03), fill_color=GOLD)
        # 描述
        for j, line in enumerate(desc.split('\n')):
            add_textbox(slide, left + Inches(0.05), top + Inches(2.4 + j * 0.4), Inches(1.6), Inches(0.35),
                        line, font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # 底部箭头连接线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.04), fill_color=YIWU_RED)
    add_textbox(slide, Inches(2.5), Inches(6.6), Inches(8), Inches(0.4),
                '7大AI Agent协同  →  从洞察到成交+政策复制一站式解决',
                font_size=16, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 5)


# ==================== 第6页：市场洞察Agent ====================
def slide_agent_market():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '市场洞察 Agent', '实时捕捉全球市场趋势，义乌指数赋能决策')

    # 左侧功能列表
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(0.4),
                '核心功能', font_size=22, font_color=YIWU_RED, bold=True)

    features = [
        '实时监控全球200+国家/地区小商品需求趋势',
        '义乌指数深度分析：价格指数、景气指数、订单指数',
        'AI预测热门品类与季节性需求波动',
        '竞品动态追踪与市场空白点识别',
        '多维度数据可视化看板，一键生成洞察报告',
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.6), Inches(5.2), Inches(3.5),
                    features, font_size=15, font_color=DARK_GRAY, bullet_char='▸')

    # 右侧义乌指数卡片
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.4),
                '义乌指数赋能', font_size=22, font_color=YIWU_RED, bold=True)

    index_data = [
        ('价格指数', '102.8', '↑ 2.3%'),
        ('景气指数', '1,256', '↑ 5.1%'),
        ('订单指数', '89.6', '↑ 3.7%'),
    ]
    for i, (name, value, change) in enumerate(index_data):
        y = Inches(2.7 + i * 1.3)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), y, Inches(5.2), Inches(1.0),
                  fill_color=LIGHT_GRAY, line_color=None)
        add_textbox(slide, Inches(7.3), y + Inches(0.1), Inches(1.8), Inches(0.35),
                    name, font_size=14, font_color=MID_GRAY)
        add_textbox(slide, Inches(9.2), y + Inches(0.05), Inches(1.8), Inches(0.5),
                    value, font_size=28, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(11.0), y + Inches(0.15), Inches(1.2), Inches(0.35),
                    change, font_size=16, font_color=ACCENT_GREEN, bold=True)

    add_page_number(slide, 6)


# ==================== 第7页：智能选品Agent ====================
def slide_agent_selection():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '智能选品 Agent', '210万SKU直连，AI驱动精准选品')

    # 左侧大数字
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(4), Inches(5.2),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(3.6), Inches(1.5),
                '210万', font_size=72, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(3.6), Inches(0.5),
                'SKU 直连', font_size=28, font_color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(3.6), Inches(1.5),
                '覆盖义乌国际商贸城\n全品类商品数据\n实时更新库存与价格',
                font_size=15, font_color=RGBColor(0xFF, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    # 右侧功能
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.8), Inches(7.6), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(5.3), Inches(2.0), Inches(7.0), Inches(0.4),
                '智能选品能力', font_size=22, font_color=YIWU_RED, bold=True)

    capabilities = [
        ('多维度筛选', '按品类、价格、销量、评分、趋势等多维度智能筛选'),
        ('AI推荐引擎', '基于目标市场偏好与历史数据的个性化推荐'),
        ('趋势预测', '预测未来30/60/90天热门品类走势'),
        ('利润测算', '自动计算成本、物流、关税后的预估利润'),
        ('一键上架', '选品后一键同步至Amazon/Shopee/Temu等平台'),
    ]
    for i, (title, desc) in enumerate(capabilities):
        y = Inches(2.6 + i * 0.85)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), y, Inches(7.0), Inches(0.7),
                  fill_color=LIGHT_GRAY, line_color=None)
        add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(2.0), Inches(0.3),
                    title, font_size=14, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(5.5), y + Inches(0.35), Inches(6.5), Inches(0.3),
                    desc, font_size=12, font_color=MID_GRAY)

    add_page_number(slide, 7)


# ==================== 第8页：供应链匹配Agent（重点页） ====================
def slide_agent_supply_chain():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '供应链匹配 Agent  ★', '义乌独有壁垒：7.5万商铺精准匹配')

    # 顶部重点标识
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.6),
              fill_color=RGBColor(0xFF, 0xF0, 0xF0), line_color=YIWU_RED, line_width=Pt(1.5))
    add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.5),
                '★ 义乌独有优势 — 全球唯一7.5万商铺实时数据直连，这是其他任何平台无法复制的核心壁垒',
                font_size=16, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)

    # 左侧核心数据
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.5),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(0.5),
                '义乌商铺直连', font_size=24, font_color=WHITE, bold=True)

    stats = [
        ('7.5万+', '国际商贸城商铺'),
        ('210万', '在线SKU商品'),
        ('实时', '库存价格更新'),
        ('100%', '商铺资质认证'),
    ]
    for i, (num, label) in enumerate(stats):
        y = Inches(3.4 + i * 0.85)
        add_textbox(slide, Inches(1.0), y, Inches(2.0), Inches(0.5),
                    num, font_size=32, font_color=GOLD, bold=True)
        add_textbox(slide, Inches(3.2), y + Inches(0.08), Inches(3.0), Inches(0.4),
                    label, font_size=16, font_color=RGBColor(0xFF, 0xDD, 0xDD))

    # 右侧匹配能力
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.5),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.5),
                '智能匹配能力', font_size=22, font_color=YIWU_RED, bold=True)

    match_features = [
        '需求画像 → 商铺画像：AI语义匹配精度95%+',
        '多商铺比价：同品类3-5家供应商智能比价',
        '产能评估：实时产能与交期预测',
        '样品快寄：一键申请样品，3天内到',
        '信用评级：商铺历史交易信用评分',
        '关系网络：老商户推荐 + 新商户发现',
    ]
    add_bullet_list(slide, Inches(7.1), Inches(3.3), Inches(5.2), Inches(3.5),
                    match_features, font_size=14, font_color=DARK_GRAY, bullet_char='▸', line_spacing=1.5)

    add_page_number(slide, 8)


# ==================== 第9页：跨境内容生成Agent ====================
def slide_agent_content():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '跨境内容生成 Agent', '8语言4平台，一键生成跨境营销内容')

    # 左侧语言和平台
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.3),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(0.9), Inches(1.95), Inches(5.2), Inches(0.4),
                '支持8种语言', font_size=20, font_color=YIWU_RED, bold=True)
    languages = ['英语', '西班牙语', '阿拉伯语', '法语', '葡萄牙语', '俄语', '日语', '韩语']
    for i, lang in enumerate(languages):
        col = i % 4
        row = i // 4
        left = Inches(0.9 + col * 1.35)
        top = Inches(2.5 + row * 0.65)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.2), Inches(0.5),
                  fill_color=RGBColor(0xFF, 0xF0, 0xF0), line_color=YIWU_RED, line_width=Pt(0.5))
        add_textbox(slide, left, top + Inches(0.05), Inches(1.2), Inches(0.35),
                    lang, font_size=13, font_color=YIWU_RED, alignment=PP_ALIGN.CENTER)

    # 平台
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.7),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(0.9), Inches(4.45), Inches(5.2), Inches(0.4),
                '覆盖4大平台', font_size=20, font_color=YIWU_RED, bold=True)
    platforms = [
        ('Amazon', 'Listing + A+内容'),
        ('Shopee', '商品描述 + 促销文案'),
        ('Temu', '商品标题 + 卖点提炼'),
        ('TikTok', '短视频脚本 + 直播话术'),
    ]
    for i, (name, desc) in enumerate(platforms):
        y = Inches(5.0 + i * 0.45)
        add_textbox(slide, Inches(1.0), y, Inches(1.5), Inches(0.35),
                    name, font_size=14, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(2.8), y, Inches(3.3), Inches(0.35),
                    desc, font_size=13, font_color=MID_GRAY)

    # 右侧内容生成能力
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.4),
                'AI内容生成能力', font_size=22, font_color=YIWU_RED, bold=True)

    content_features = [
        '商品标题：多语言SEO优化标题自动生成',
        '商品描述：符合平台规范的高转化描述',
        '营销文案：社交媒体帖子、广告语、促销文案',
        '视频脚本：TikTok/Reels短视频脚本',
        '图片文案：主图卖点标注 + 详情页文案',
        '本地化适配：文化禁忌检测 + 表达习惯优化',
    ]
    add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(4.0),
                    content_features, font_size=14, font_color=DARK_GRAY, bullet_char='▸', line_spacing=1.5)

    add_page_number(slide, 9)


# ==================== 第10页：合规助手Agent ====================
def slide_agent_compliance():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '合规助手 Agent', '1039模式 + 义新欧沿线合规，为出海保驾护航')

    # 左侧1039模式
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.85), Inches(5.7), Inches(0.08), fill_color=YIWU_RED)
    add_textbox(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.4),
                '市场采购贸易1039模式', font_size=20, font_color=YIWU_RED, bold=True)

    features_1039 = [
        '增值税免征不退，降低出口成本15-20%',
        '简化报关流程，通关时效提升50%',
        '多品类拼箱出口，小单也能享受政策红利',
        '外汇便利化，合法合规收结汇',
        '义乌试点经验，全国推广标杆',
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.7), Inches(5.2), Inches(3.5),
                    features_1039, font_size=14, font_color=DARK_GRAY, bullet_char='✓', line_spacing=1.6)

    # 右侧义新欧合规
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.85), Inches(5.7), Inches(0.08), fill_color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.4),
                '义新欧沿线合规', font_size=20, font_color=ACCENT_BLUE, bold=True)

    compliance_items = [
        '欧盟CE认证、REACH法规自动检测',
        '中亚五国贸易合规要求实时更新',
        '中东地区Halal认证指导',
        '东南亚各国进口准入标准查询',
        '沿线国家关税税率智能计算',
        '合规风险预警与整改建议',
    ]
    add_bullet_list(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(3.5),
                    compliance_items, font_size=14, font_color=DARK_GRAY, bullet_char='✓', line_spacing=1.5)

    add_page_number(slide, 10)


# ==================== 第11页：智能客服Agent ====================
def slide_agent_service():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '智能客服 Agent', '7×24小时多语言智能客服，让每个客户都被善待')

    # 中心大数字
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(4.3), Inches(2.5),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(4.7), Inches(2.0), Inches(3.9), Inches(1.0),
                '7×24', font_size=64, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.7), Inches(3.0), Inches(3.9), Inches(0.5),
                '全天候多语言服务', font_size=20, font_color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

    # 四个能力卡片
    service_features = [
        ('多语言对话', '支持8种语言实时对话\n自动识别客户语言\n文化适配表达方式'),
        ('智能工单', '自动创建跟进工单\n优先级智能排序\n超时自动升级'),
        ('知识库', '产品FAQ自动应答\n订单状态实时查询\n退换货流程指引'),
        ('人机协作', '复杂问题转人工\n对话上下文无缝传递\n服务质量实时监控'),
    ]

    for i, (title, desc) in enumerate(service_features):
        left = Inches(0.6 + i * 3.15)
        top = Inches(4.7)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(2.3),
                  fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(2.8), Inches(0.06), fill_color=YIWU_RED)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.25), Inches(2.6), Inches(0.35),
                    title, font_size=17, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)
        for j, line in enumerate(desc.split('\n')):
            add_textbox(slide, left + Inches(0.15), top + Inches(0.75 + j * 0.4), Inches(2.6), Inches(0.35),
                        line, font_size=12, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 11)


# ==================== 第12页：政策复制Agent（新增） ====================
def slide_agent_policy():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '政策复制 Agent  ★', '第7大Agent：义乌经验39城一键复制')

    # 顶部标识
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.6),
              fill_color=RGBColor(0xFF, 0xF8, 0xE0), line_color=GOLD, line_width=Pt(1.5))
    add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.5),
                '★ 国家级政策复制能力 — 1039模式已在全国39个城市推广，OPC是唯一AI驱动的政策复制引擎',
                font_size=16, font_color=RGBColor(0x99, 0x66, 0x00), bold=True, alignment=PP_ALIGN.CENTER)

    # 左侧核心能力
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.5),
              fill_color=GOLD, line_color=None)
    add_textbox(slide, Inches(0.9), Inches(2.7), Inches(5.2), Inches(0.5),
                '政策复制核心能力', font_size=22, font_color=DARK_RED, bold=True)

    policy_features = [
        ('政策匹配', 'AI分析目标城市产业特征，匹配最优义乌经验方案'),
        ('流程复制', '1039通关流程、市场采购模式一键复制到新城市'),
        ('合规迁移', '各地法规差异自动识别，合规方案本地化适配'),
        ('数据互通', '39城政策数据网络互通，形成全国政策知识图谱'),
        ('效果追踪', '复制效果实时监控，AI持续优化落地策略'),
    ]
    for i, (title, desc) in enumerate(policy_features):
        y = Inches(3.4 + i * 0.7)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(5.2), Inches(0.6),
                  fill_color=RGBColor(0xC8, 0x96, 0x10), line_color=None)
        add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(1.5), Inches(0.3),
                    title, font_size=14, font_color=WHITE, bold=True)
        add_textbox(slide, Inches(2.7), y + Inches(0.05), Inches(3.2), Inches(0.5),
                    desc, font_size=12, font_color=RGBColor(0x4A, 0x2A, 0x00))

    # 右侧39城数据
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.5),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.5),
                '39城复制推广数据', font_size=22, font_color=YIWU_RED, bold=True)

    city_stats = [
        ('39', '试点城市', '覆盖全国主要贸易节点城市'),
        ('210万', '市场主体', '义乌经验带动全国市场主体发展'),
        ('3200万', '就业人口', '义乌模式支撑就业规模'),
        ('50万', '政府年采购', '每城市政府年采购额（元）'),
    ]
    for i, (num, label, desc) in enumerate(city_stats):
        y = Inches(3.4 + i * 0.85)
        add_textbox(slide, Inches(7.3), y, Inches(1.5), Inches(0.5),
                    num, font_size=28, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(8.9), y, Inches(1.5), Inches(0.3),
                    label, font_size=14, font_color=DARK_GRAY, bold=True)
        add_textbox(slide, Inches(8.9), y + Inches(0.3), Inches(3.2), Inches(0.3),
                    desc, font_size=11, font_color=MID_GRAY)

    add_page_number(slide, 12)


# ==================== 第13页：39城复制推广地图页（新增） ====================
def slide_city_map():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '39城复制推广地图', '1039模式全国布局，从义乌到全国')

    # 中心地图区域
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(8.5), Inches(5.2),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))

    # 义乌中心标识
    add_shape(slide, MSO_SHAPE.OVAL, Inches(4.2), Inches(3.5), Inches(1.2), Inches(1.2),
              fill_color=YIWU_RED, line_color=GOLD, line_width=Pt(3))
    add_textbox(slide, Inches(4.2), Inches(3.75), Inches(1.2), Inches(0.7),
                '义乌', font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 城市分布点（模拟地图布局）
    cities = [
        # (x_offset, y_offset, name, region)
        (1.5, 2.2, '满洲里', '东北'),
        (2.5, 2.5, '绥芬河', '东北'),
        (3.5, 2.0, '黑河', '东北'),
        (5.5, 2.2, '丹东', '东北'),
        (2.0, 3.0, '二连浩特', '华北'),
        (3.0, 3.2, '呼和浩特', '华北'),
        (4.0, 3.0, '北京', '华北'),
        (5.0, 3.5, '天津', '华北'),
        (3.5, 3.8, '石家庄', '华北'),
        (4.5, 3.8, '济南', '华东'),
        (5.5, 3.8, '青岛', '华东'),
        (6.5, 3.5, '上海', '华东'),
        (6.0, 4.2, '杭州', '华东'),
        (5.0, 4.5, '南昌', '华中'),
        (4.5, 4.8, '长沙', '华中'),
        (3.5, 4.5, '武汉', '华中'),
        (4.0, 5.2, '广州', '华南'),
        (4.5, 5.5, '深圳', '华南'),
        (5.0, 5.2, '厦门', '华南'),
        (3.0, 5.0, '南宁', '华南'),
        (2.0, 4.5, '昆明', '西南'),
        (2.5, 4.0, '成都', '西南'),
        (2.0, 3.8, '重庆', '西南'),
        (1.5, 3.5, '西安', '西北'),
        (1.0, 3.0, '兰州', '西北'),
        (0.8, 2.5, '乌鲁木齐', '西北'),
        (3.0, 4.2, '贵阳', '西南'),
        (3.5, 5.0, '海口', '华南'),
        (6.5, 4.0, '宁波', '华东'),
        (5.5, 4.5, '福州', '华东'),
        (6.0, 5.0, '泉州', '华东'),
        (3.0, 3.5, '太原', '华北'),
        (2.5, 3.5, '银川', '西北'),
        (1.5, 4.0, '拉萨', '西南'),
        (4.0, 4.2, '合肥', '华东'),
        (5.0, 4.0, '南京', '华东'),
        (4.5, 3.5, '郑州', '华中'),
        (3.0, 4.8, '凭祥', '华南'),
        (2.0, 2.8, '阿拉山口', '西北'),
    ]

    for i, (x, y, name, region) in enumerate(cities):
        left = Inches(0.8 + x * 0.9)
        top = Inches(1.9 + y * 0.55)
        # 城市点
        dot_size = Inches(0.2)
        add_shape(slide, MSO_SHAPE.OVAL, left, top, dot_size, dot_size,
                  fill_color=ACCENT_BLUE if region != '华东' else YIWU_RED)
        # 城市名
        add_textbox(slide, left - Inches(0.15), top + Inches(0.18), Inches(0.6), Inches(0.25),
                    name, font_size=7, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # 右侧统计面板
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(1.8), Inches(3.2), Inches(5.2),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(9.7), Inches(2.0), Inches(2.8), Inches(0.5),
                '39城布局概览', font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    region_stats = [
        ('华东', '12城', '义乌、上海、杭州、宁波等'),
        ('华南', '7城', '广州、深圳、厦门、南宁等'),
        ('华北', '6城', '北京、天津、石家庄等'),
        ('华中', '5城', '武汉、长沙、郑州等'),
        ('东北', '4城', '满洲里、绥芬河、丹东等'),
        ('西南', '3城', '成都、重庆、昆明等'),
        ('西北', '2城', '西安、乌鲁木齐等'),
    ]
    for i, (region, count, cities_str) in enumerate(region_stats):
        y = Inches(2.7 + i * 0.6)
        add_textbox(slide, Inches(9.7), y, Inches(1.0), Inches(0.25),
                    region, font_size=12, font_color=GOLD, bold=True)
        add_textbox(slide, Inches(10.7), y, Inches(0.8), Inches(0.25),
                    count, font_size=12, font_color=WHITE, bold=True)
        add_textbox(slide, Inches(9.7), y + Inches(0.25), Inches(2.8), Inches(0.25),
                    cities_str, font_size=9, font_color=RGBColor(0xFF, 0xCC, 0xCC))

    add_page_number(slide, 13)


# ==================== 第14页：全链路工作流（7步） ====================
def slide_workflow():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '全链路工作流', '7步串联，从市场洞察到政策复制全闭环')

    steps = [
        ('01', '市场洞察', '发现机会', '趋势分析\n需求预测'),
        ('02', '智能选品', '锁定产品', 'SKU筛选\n利润测算'),
        ('03', '供应链匹配', '对接商铺', '商铺匹配\n样品确认'),
        ('04', '内容生成', '打造内容', '多语言文案\n平台适配'),
        ('05', '合规审核', '确保合规', '资质检测\n风险预警'),
        ('06', '智能客服', '成交转化', '客户服务\n订单跟进'),
        ('07', '政策复制', '全国推广', '39城复制\n经验迁移'),
    ]

    for i, (num, title, slogan, desc) in enumerate(steps):
        left = Inches(0.3 + i * 1.85)
        top = Inches(2.0)

        # 步骤卡片
        color = YIWU_RED if i < 6 else GOLD
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.7), Inches(4.5),
                  fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        # 顶部编号
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(1.6), Inches(0.7), fill_color=color)
        add_textbox(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.6), Inches(0.6),
                    num, font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # 步骤标题
        add_textbox(slide, left + Inches(0.05), top + Inches(0.95), Inches(1.6), Inches(0.4),
                    title, font_size=15, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)
        # Slogan
        add_textbox(slide, left + Inches(0.05), top + Inches(1.4), Inches(1.6), Inches(0.3),
                    slogan, font_size=12, font_color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
        # 分隔线
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.3), top + Inches(1.85),
                  Inches(1.1), Inches(0.03), fill_color=RGBColor(0xDD, 0xDD, 0xDD))
        # 描述
        for j, line in enumerate(desc.split('\n')):
            add_textbox(slide, left + Inches(0.1), top + Inches(2.1 + j * 0.4), Inches(1.5), Inches(0.35),
                        line, font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # 箭头（最后一个不加）
        if i < 6:
            arrow_left = left + Inches(1.7)
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, arrow_left, top + Inches(2.0),
                      Inches(0.15), Inches(0.3), fill_color=YIWU_RED)

    # 底部说明
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.5),
              fill_color=RGBColor(0xFF, 0xF0, 0xF0), line_color=YIWU_RED, line_width=Pt(0.5))
    add_textbox(slide, Inches(1.7), Inches(6.75), Inches(9.9), Inches(0.4),
                '7大Agent协同工作，数据自动流转，从义乌到39城全国复制',
                font_size=15, font_color=YIWU_RED, bold=True, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 14)


# ==================== 第15页：技术架构页 ====================
def slide_tech_arch():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '技术架构', 'LangGraph + FastAPI + React，企业级AI架构')

    layers = [
        ('前端层', 'React + Next.js', [
            '响应式Web应用',
            '实时数据看板',
            '多语言UI支持',
            'PWA离线可用',
        ], ACCENT_BLUE),
        ('服务层', 'FastAPI + LangGraph', [
            'LangGraph多Agent编排',
            'FastAPI高性能接口',
            'Redis消息队列',
            'PostgreSQL数据持久化',
        ], YIWU_RED),
        ('数据层', '多源数据融合', [
            '义乌指数API',
            '商铺实时数据',
            '海关合规数据库',
            '39城政策数据库',
        ], ACCENT_GREEN),
    ]

    for i, (layer_name, tech, items, color) in enumerate(layers):
        left = Inches(0.6 + i * 4.2)
        top = Inches(1.8)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(0.7),
                  fill_color=color, line_color=None)
        add_textbox(slide, left, top + Inches(0.1), Inches(3.8), Inches(0.5),
                    f'{layer_name}  |  {tech}', font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.8), Inches(3.8), Inches(3.5),
                  fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        add_bullet_list(slide, left + Inches(0.3), top + Inches(1.1), Inches(3.2), Inches(3.0),
                        items, font_size=14, font_color=DARK_GRAY, bullet_char='▸', line_spacing=1.8)

    # 底部AI能力
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0),
              fill_color=LIGHT_GRAY, line_color=None)
    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
                'AI核心能力', font_size=16, font_color=YIWU_RED, bold=True)
    ai_caps = 'GPT-4o 多语言理解  |  LangGraph Agent编排  |  RAG知识增强  |  向量检索  |  实时数据流  |  多模态内容生成  |  政策知识图谱'
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.35),
                ai_caps, font_size=13, font_color=MID_GRAY)

    add_page_number(slide, 15)


# ==================== 第16页：商业模式页 ====================
def slide_business_model():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '商业模式', '四档定价 + 政府采购 + 城市合伙人，多元收入')

    # 定价表
    plans = [
        ('体验版', '免费', '0', [
            '每月50次AI调用',
            '基础市场数据',
            '1个平台对接',
            '社区支持',
        ], RGBColor(0x99, 0x99, 0x99)),
        ('基础版', '¥299/月', '299', [
            '每月500次AI调用',
            '完整市场洞察',
            '3个平台对接',
            '在线客服支持',
        ], ACCENT_BLUE),
        ('专业版', '¥999/月', '999', [
            '无限AI调用',
            '全链路7大Agent',
            '全平台对接',
            '专属客户经理',
        ], YIWU_RED),
        ('企业版', '定制报价', '-', [
            '私有化部署',
            '定制Agent开发',
            'API接口开放',
            '7×24专属服务',
        ], ACCENT_GREEN),
    ]

    for i, (name, price, price_num, features, color) in enumerate(plans):
        left = Inches(0.5 + i * 3.2)
        top = Inches(1.8)

        is_pro = (name == '专业版')
        card_h = Inches(5.5) if is_pro else Inches(5.0)
        card_top = Inches(1.8) if is_pro else Inches(2.1)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, Inches(2.9), card_h,
                  fill_color=WHITE, line_color=color, line_width=Pt(2 if is_pro else 1))

        if is_pro:
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.7), card_top - Inches(0.25),
                      Inches(1.5), Inches(0.4), fill_color=YIWU_RED, line_color=None)
            add_textbox(slide, left + Inches(0.7), card_top - Inches(0.22), Inches(1.5), Inches(0.35),
                        '推荐', font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), card_top + Inches(0.05),
                  Inches(2.8), Inches(0.8), fill_color=color)
        add_textbox(slide, left + Inches(0.1), card_top + Inches(0.1), Inches(2.7), Inches(0.35),
                    name, font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), card_top + Inches(0.45), Inches(2.7), Inches(0.35),
                    price, font_size=16, font_color=RGBColor(0xFF, 0xEE, 0xEE), alignment=PP_ALIGN.CENTER)

        for j, feat in enumerate(features):
            add_textbox(slide, left + Inches(0.2), card_top + Inches(1.1 + j * 0.55), Inches(2.5), Inches(0.4),
                        f'✓  {feat}', font_size=13, font_color=DARK_GRAY)

    add_page_number(slide, 16)


# ==================== 第17页：政府采购+城市合伙人商业模式（新增） ====================
def slide_gov_partner_model():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '政府采购 + 城市合伙人', 'B2G + B2B双轮驱动，39城规模化复制')

    # 左侧：政府采购
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=YIWU_RED, line_width=Pt(2))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.85), Inches(5.7), Inches(0.8),
              fill_color=YIWU_RED)
    add_textbox(slide, Inches(0.8), Inches(1.95), Inches(5.4), Inches(0.6),
                '政府采购模式（B2G）', font_size=22, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    gov_items = [
        ('年费', '50万/年/城市', '政府购买义乌经验数字化复制服务'),
        ('覆盖', '39个试点城市', '1039模式推广城市全覆盖'),
        ('内容', '政策复制+合规+数据', '一站式数字化出海公共服务'),
        ('收入', '1,950万/年', '39城 × 50万 = 1,950万/年'),
    ]
    for i, (label, value, desc) in enumerate(gov_items):
        y = Inches(2.9 + i * 0.95)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(5.2), Inches(0.8),
                  fill_color=LIGHT_GRAY, line_color=None)
        add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(1.2), Inches(0.3),
                    label, font_size=14, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(2.4), y + Inches(0.05), Inches(3.5), Inches(0.3),
                    value, font_size=16, font_color=DARK_GRAY, bold=True)
        add_textbox(slide, Inches(1.1), y + Inches(0.4), Inches(4.8), Inches(0.3),
                    desc, font_size=12, font_color=MID_GRAY)

    # 右侧：城市合伙人
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2),
              fill_color=WHITE, line_color=GOLD, line_width=Pt(2))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.85), Inches(5.7), Inches(0.8),
              fill_color=GOLD)
    add_textbox(slide, Inches(7.0), Inches(1.95), Inches(5.4), Inches(0.6),
                '城市合伙人模式（B2B）', font_size=22, font_color=DARK_RED, bold=True, alignment=PP_ALIGN.CENTER)

    partner_items = [
        ('分润比例', '30%', '合伙人获得当地商户收入的30%'),
        ('合伙人职责', '本地运营+商户拓展', '负责本地市场推广与商户服务'),
        ('OPC职责', '技术+品牌+培训', '提供平台、培训、品牌支持'),
        ('预期收入', '585万/年', '39城商户收入的30%分润'),
    ]
    for i, (label, value, desc) in enumerate(partner_items):
        y = Inches(2.9 + i * 0.95)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), y, Inches(5.2), Inches(0.8),
                  fill_color=RGBColor(0xFF, 0xF8, 0xE0), line_color=None)
        add_textbox(slide, Inches(7.3), y + Inches(0.05), Inches(1.5), Inches(0.3),
                    label, font_size=14, font_color=RGBColor(0x99, 0x66, 0x00), bold=True)
        add_textbox(slide, Inches(8.9), y + Inches(0.05), Inches(3.2), Inches(0.3),
                    value, font_size=16, font_color=DARK_GRAY, bold=True)
        add_textbox(slide, Inches(7.3), y + Inches(0.4), Inches(4.8), Inches(0.3),
                    desc, font_size=12, font_color=MID_GRAY)

    add_page_number(slide, 17)


# ==================== 第18页：竞争优势页（增加义乌经验壁垒） ====================
def slide_competitive():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '竞争优势', '四大义乌独有壁垒，构建不可逾越的护城河')

    advantages = [
        ('壁垒一', '7.5万商铺数据直连', YIWU_RED,
         '全球唯一与义乌国际商贸城7.5万商铺实时数据直连的平台',
         ['商铺库存、价格实时同步', '商户信用评级体系', '供应链响应速度领先3-5倍', '数据壁垒：10年+积累']),
        ('壁垒二', '1039政策深度绑定', ACCENT_BLUE,
         '深度绑定义乌市场采购贸易1039模式，享受政策红利',
         ['增值税免征不退，成本优势15-20%', '简化通关，时效提升50%', '政府背书，合规无忧', '政策壁垒：试点独占']),
        ('壁垒三', '义新欧通道优势', ACCENT_GREEN,
         '依托义新欧中欧班列，构建沿线合规与服务网络',
         ['覆盖中亚、欧洲13国', '沿线合规数据库独家建设', '物流时效提升30%', '通道壁垒：先发优势']),
        ('壁垒四', '义乌经验国家级壁垒', GOLD,
         '习近平总书记亲自批示"义乌发展经验"，39城政策复制',
         ['国家级战略背书，不可复制', '39城政策网络，规模壁垒', '210万企业+3200万就业数据飞轮', '制度壁垒：唯一性']),
    ]

    for i, (tag, title, color, desc, points) in enumerate(advantages):
        left = Inches(0.3 + i * 3.25)
        top = Inches(1.8)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.0), Inches(5.2),
                  fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        # 顶部色块
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(2.9), Inches(1.0), fill_color=color)
        # 标签
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(1.5), Inches(0.3),
                    tag, font_size=12, font_color=RGBColor(0xFF, 0xEE, 0xEE))
        # 标题
        add_textbox(slide, left + Inches(0.15), top + Inches(0.45), Inches(2.7), Inches(0.5),
                    title, font_size=16, font_color=WHITE, bold=True)
        # 描述
        add_textbox(slide, left + Inches(0.15), top + Inches(1.3), Inches(2.7), Inches(0.8),
                    desc, font_size=11, font_color=MID_GRAY)
        # 分隔线
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(2.2),
                  Inches(2.6), Inches(0.03), fill_color=color)
        # 要点
        add_bullet_list(slide, left + Inches(0.15), top + Inches(2.5), Inches(2.7), Inches(2.5),
                        points, font_size=11, font_color=DARK_GRAY, bullet_char='▸', line_spacing=1.5)

    add_page_number(slide, 18)


# ==================== 第19页：OPC模式页 ====================
def slide_opc_model():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, 'OPC 模式', '1人 + 7Agent = 10人团队效能')

    # 中心公式
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.0),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(1.7), Inches(2.0), Inches(9.9), Inches(0.8),
                '1人 + 7Agent = 10人团队', font_size=44, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.7), Inches(2.9), Inches(9.9), Inches(0.5),
                'One Person Company — 一个人就是一家出海公司', font_size=20, font_color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

    # 对比表
    # 传统模式
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.8),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(4.25), Inches(5.7), Inches(0.6),
              fill_color=MID_GRAY)
    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.5),
                '传统模式', font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    trad_items = ['需要10+人团队', '月成本10万+', '响应周期3-7天', '人工操作易出错']
    for i, item in enumerate(trad_items):
        add_textbox(slide, Inches(1.0), Inches(5.1 + i * 0.45), Inches(5.0), Inches(0.35),
                    f'✗  {item}', font_size=14, font_color=RGBColor(0xCC, 0x33, 0x33))

    # OPC模式
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.8),
              fill_color=WHITE, line_color=YIWU_RED, line_width=Pt(2))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(4.25), Inches(5.7), Inches(0.6),
              fill_color=YIWU_RED)
    add_textbox(slide, Inches(7.0), Inches(4.3), Inches(5.4), Inches(0.5),
                'OPC模式', font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    opc_items = ['1人即可运营', '月成本低至999元', 'AI实时响应，秒级决策', '7大Agent自动化，零出错']
    for i, item in enumerate(opc_items):
        add_textbox(slide, Inches(7.2), Inches(5.1 + i * 0.45), Inches(5.0), Inches(0.35),
                    f'✓  {item}', font_size=14, font_color=ACCENT_GREEN, bold=True)

    add_page_number(slide, 19)


# ==================== 第20页：三级增长飞轮页（新增） ====================
def slide_growth_flywheel():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '三级增长飞轮', '义乌1.0 → 39城2.0 → 全国3.0，指数级增长')

    # 三个飞轮层级
    flywheels = [
        ('第一级', '义乌本地 1.0', '1-12个月', YIWU_RED,
         ['7.5万商铺直连', '210万SKU数据', '1039模式试点', 'MRR 144.7万'],
         '夯实义乌根据地，验证商业模式'),
        ('第二级', '39城复制 2.0', '13-24个月', GOLD,
         ['39城政策复制', '政府采购1,950万/年', '城市合伙人30%分润', '商户规模10倍增长'],
         '义乌经验全国复制，B2G+B2B双轮驱动'),
        ('第三级', '全国县域 3.0', '25-36个月', ACCENT_BLUE,
         ['2,800+县域覆盖', '全国小商品网络', '产业带数字化', '平台生态收入'],
         '从城市到县域，构建全国小商品数字化基础设施'),
    ]

    for i, (level, title, period, color, points, slogan) in enumerate(flywheels):
        left = Inches(0.5 + i * 4.2)
        top = Inches(1.8)

        # 层级卡片
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.9), Inches(5.2),
                  fill_color=WHITE, line_color=color, line_width=Pt(2))

        # 顶部色块
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                  Inches(3.8), Inches(1.2), fill_color=color)
        # 层级标签
        add_textbox(slide, left + Inches(0.2), top + Inches(0.1), Inches(1.5), Inches(0.3),
                    level, font_size=13, font_color=RGBColor(0xFF, 0xEE, 0xEE))
        # 标题
        add_textbox(slide, left + Inches(0.2), top + Inches(0.45), Inches(3.5), Inches(0.5),
                    title, font_size=22, font_color=WHITE, bold=True)
        # 周期
        add_textbox(slide, left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(0.3),
                    period, font_size=14, font_color=RGBColor(0xFF, 0xDD, 0xDD))

        # 要点
        for j, point in enumerate(points):
            y = top + Inches(1.5 + j * 0.5)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), y, Inches(3.5), Inches(0.4),
                      fill_color=LIGHT_GRAY, line_color=None)
            add_textbox(slide, left + Inches(0.35), y + Inches(0.05), Inches(3.2), Inches(0.3),
                        f'▸  {point}', font_size=13, font_color=DARK_GRAY, bold=False)

        # 底部Slogan
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.3), top + Inches(3.7),
                  Inches(3.3), Inches(0.03), fill_color=color)
        add_textbox(slide, left + Inches(0.2), top + Inches(3.9), Inches(3.5), Inches(1.0),
                    slogan, font_size=12, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # 飞轮箭头
        if i < 2:
            arrow_left = left + Inches(3.9)
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, arrow_left, top + Inches(2.5),
                      Inches(0.3), Inches(0.5), fill_color=color)

    add_page_number(slide, 20)


# ==================== 第21页：财务预测页（含39城收入） ====================
def slide_financial():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '财务预测', '12个月收入+39城复制，第8个月盈亏平衡')

    # 收入数据表
    months_data = [
        ('月份', '付费用户', '月收入(万)', '累计收入(万)', '月成本(万)', '净利润(万)'),
        ('M1-3', '50', '3.0', '3.0', '5.0', '-2.0'),
        ('M4-6', '200', '15.0', '18.0', '6.0', '9.0'),
        ('M7-9', '500', '40.0', '58.0', '8.0', '32.0'),
        ('M10-12', '1000', '85.0', '143.0', '12.0', '73.0'),
    ]

    # 表格
    table_left = Inches(0.6)
    table_top = Inches(1.8)
    col_widths = [Inches(1.5), Inches(2.0), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.0)]

    for row_idx, row_data in enumerate(months_data):
        for col_idx, cell_text in enumerate(row_data):
            left = table_left
            for c in range(col_idx):
                left += col_widths[c]

            is_header = (row_idx == 0)
            bg = YIWU_RED if is_header else (LIGHT_GRAY if row_idx % 2 == 0 else WHITE)
            fc = WHITE if is_header else DARK_GRAY
            row_height = Inches(0.5)

            add_shape(slide, MSO_SHAPE.RECTANGLE, left, table_top + row_idx * row_height,
                      col_widths[col_idx], row_height, fill_color=bg,
                      line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
            add_textbox(slide, left + Inches(0.1), table_top + row_idx * row_height + Inches(0.08),
                        col_widths[col_idx] - Inches(0.2), row_height - Inches(0.16),
                        cell_text, font_size=13, font_color=fc, bold=is_header, alignment=PP_ALIGN.CENTER)

    # 39城收入预测
    add_textbox(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
                '39城复制收入预测（M13-24）', font_size=18, font_color=YIWU_RED, bold=True)

    city_revenue_data = [
        ('项目', 'M13-15', 'M16-18', 'M19-21', 'M22-24'),
        ('政府采购(万/季)', '375', '750', '1,125', '1,500'),
        ('城市合伙人分润(万/季)', '113', '225', '338', '450'),
        ('39城商户收入(万/季)', '263', '525', '788', '1,050'),
        ('合计(万/季)', '751', '1,500', '2,251', '3,000'),
    ]

    city_table_top = Inches(4.7)
    city_col_widths = [Inches(3.0), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]

    for row_idx, row_data in enumerate(city_revenue_data):
        for col_idx, cell_text in enumerate(row_data):
            left = Inches(0.6)
            for c in range(col_idx):
                left += city_col_widths[c]

            is_header = (row_idx == 0)
            is_total = (row_idx == 4)
            bg = YIWU_RED if is_header else (GOLD if is_total else (LIGHT_GRAY if row_idx % 2 == 0 else WHITE))
            fc = WHITE if is_header or is_total else DARK_GRAY
            row_height = Inches(0.4)

            add_shape(slide, MSO_SHAPE.RECTANGLE, left, city_table_top + row_idx * row_height,
                      city_col_widths[col_idx], row_height, fill_color=bg,
                      line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
            add_textbox(slide, left + Inches(0.1), city_table_top + row_idx * row_height + Inches(0.05),
                        city_col_widths[col_idx] - Inches(0.2), row_height - Inches(0.1),
                        cell_text, font_size=12, font_color=fc, bold=is_header or is_total, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 21)


# ==================== 第22页：融资计划页 ====================
def slide_funding():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, '融资计划', '种子轮100万，加速产品落地与39城复制')

    # 融资金额
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.5),
              fill_color=YIWU_RED, line_color=None)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.4), Inches(0.5),
                '种子轮融资', font_size=22, font_color=RGBColor(0xFF, 0xDD, 0xDD))
    add_textbox(slide, Inches(0.8), Inches(2.6), Inches(5.4), Inches(1.0),
                '¥100万', font_size=60, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5.4), Inches(0.4),
                '出让10%股权  |  估值1000万', font_size=16, font_color=RGBColor(0xFF, 0xCC, 0xCC))

    # 资金用途
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.5),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.4),
                '资金用途', font_size=20, font_color=YIWU_RED, bold=True)

    uses = [
        ('35%', '产品研发', '7大Agent完善 + 政策复制Agent'),
        ('25%', '市场推广', '义乌商户拓展 + 39城复制'),
        ('25%', '团队建设', '核心岗位招聘 + 城市合伙人'),
        ('15%', '运营储备', '服务器 + 日常运营'),
    ]
    for i, (pct, name, desc) in enumerate(uses):
        y = Inches(2.5 + i * 0.42)
        add_textbox(slide, Inches(7.1), y, Inches(0.8), Inches(0.35),
                    pct, font_size=14, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(8.0), y, Inches(1.5), Inches(0.35),
                    name, font_size=14, font_color=DARK_GRAY, bold=True)
        add_textbox(slide, Inches(9.6), y, Inches(2.8), Inches(0.35),
                    desc, font_size=12, font_color=MID_GRAY)

    # 里程碑
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5),
              fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
    add_textbox(slide, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.4),
                '融资后里程碑', font_size=20, font_color=YIWU_RED, bold=True)

    milestones = [
        ('M1-3', '产品MVP上线，7大Agent核心功能完成', '义乌本地种子用户50+'),
        ('M4-6', '付费用户200+，月收入15万', '启动39城政策复制对接'),
        ('M7-9', '盈亏平衡，付费用户500+', '首批5城政策复制落地'),
        ('M10-12', '月收入85万，准备Pre-A轮', '39城全面复制启动'),
    ]
    for i, (period, goal, detail) in enumerate(milestones):
        y = Inches(5.4 + i * 0.42)
        add_shape(slide, MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05), Inches(0.25), Inches(0.25),
                  fill_color=YIWU_RED)
        add_textbox(slide, Inches(1.4), y, Inches(1.5), Inches(0.35),
                    period, font_size=13, font_color=YIWU_RED, bold=True)
        add_textbox(slide, Inches(3.0), y, Inches(5.0), Inches(0.35),
                    goal, font_size=13, font_color=DARK_GRAY, bold=True)
        add_textbox(slide, Inches(8.2), y, Inches(4.0), Inches(0.35),
                    detail, font_size=12, font_color=MID_GRAY)

    add_page_number(slide, 22)


# ==================== 主函数 ====================
def main():
    print('开始生成路演PPT V2冠军版...')

    # 生成所有页面
    slide_cover()               # 1. 封面页
    slide_pain_points()         # 2. 痛点页
    slide_market_size()         # 3. 市场规模页
    slide_yiwu_strategy()       # 4. 义乌发展经验国家战略页（新增）
    slide_solution_overview()   # 5. 解决方案总览页（7个Agent）
    slide_agent_market()        # 6. 市场洞察Agent
    slide_agent_selection()     # 7. 智能选品Agent
    slide_agent_supply_chain()  # 8. 供应链匹配Agent（重点）
    slide_agent_content()       # 9. 跨境内容生成Agent
    slide_agent_compliance()    # 10. 合规助手Agent
    slide_agent_service()       # 11. 智能客服Agent
    slide_agent_policy()        # 12. 政策复制Agent（新增）
    slide_city_map()            # 13. 39城复制推广地图页（新增）
    slide_workflow()            # 14. 全链路工作流（7步）
    slide_tech_arch()           # 15. 技术架构
    slide_business_model()      # 16. 商业模式
    slide_gov_partner_model()   # 17. 政府采购+城市合伙人商业模式（新增）
    slide_competitive()         # 18. 竞争优势（4大壁垒）
    slide_opc_model()           # 19. OPC模式
    slide_growth_flywheel()     # 20. 三级增长飞轮（新增）
    slide_financial()           # 21. 财务预测（含39城收入）
    slide_funding()             # 22. 融资计划

    # 保存文件
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '路演PPT_V2.pptx')
    prs.save(output_path)
    print(f'PPT V2冠军版生成成功！文件路径：{output_path}')
    print(f'共 {len(prs.slides)} 页')


if __name__ == '__main__':
    main()
